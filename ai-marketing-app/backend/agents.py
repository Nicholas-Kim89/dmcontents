"""
Multi-Agent B2B Marketing Pipeline using LangGraph + LangChain Google GenAI.

Agents:
  1. Researcher: B2B Market & Competitor Intelligence
  2. Planner:    Strategy Playbook (ROI, Compliance, Specs, etc.)
  3. Creator:    Distinct B2B target copy variants (Value, Tech, Trust)
  4. Designer:   B2B professional visual image prompts
  5. Critic:     Factual & Brand safety audit (APPROVE or REVISE)
"""

import os
import io
import base64
import logging
from typing import TypedDict, Annotated, Optional, List
from typing import Literal
import operator

from langchain_core.messages import HumanMessage, SystemMessage
from langgraph.graph import StateGraph, END
from google import genai
from google.genai import types

logger = logging.getLogger("agents")

# ---------------------------------------------------------------------------
# 1. State definition
# ---------------------------------------------------------------------------

class CampaignState(TypedDict):
    # Input
    user_prompt: str
    knowledge_text: str          # extracted text from uploaded files
    target_industry: str         # B2B industry focus
    buyer_persona: str           # B2B target persona/role
    b2b_strategy: str            # B2B campaign goal
    graph_id: str                # Unique ID for progress tracking

    # Agent outputs
    research_data: str           # Researcher output (market intelligence)
    competitor_data: str         # Competitor agent output (competitor analysis)
    keyword_data: str            # Keyword agent output (keyword analysis)
    strategy: str                # Planner output (strategy playbook)
    copies: List[str]            # Creator output – list of B2B copy variants
    image_prompts: List[str]     # Designer output – premium image generation prompts
    review_result: str           # Critic output: "APPROVE" or "REVISE: <reason>"

    # Control
    retry_count: int
    final_copies: List[str]
    final_image_prompts: List[str]
    error: Optional[str]
    google_search: bool

# ---------------------------------------------------------------------------
# Progress Callback Registry (used by main.py to track real-time agent stage)
# ---------------------------------------------------------------------------

_progress_callbacks: dict = {}

def register_progress_callback(graph_id: str, callback) -> None:
    """Register a callback fn(stage: str) keyed by graph_id."""
    _progress_callbacks[graph_id] = callback

def unregister_progress_callback(graph_id: str) -> None:
    """Clean up after pipeline finishes."""
    _progress_callbacks.pop(graph_id, None)

def _report_stage(state: dict, stage_name: str) -> None:
    """Call the registered progress callback for this graph run."""
    graph_id = state.get("graph_id", "")
    cb = _progress_callbacks.get(graph_id)
    if cb:
        try:
            cb(stage_name)
        except Exception:
            pass

def _get_client(api_key: str) -> genai.Client:
    return genai.Client(api_key=api_key)

def _extract_text(content) -> str:
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        texts = []
        for part in content:
            if isinstance(part, str):
                texts.append(part)
            elif isinstance(part, dict) and "text" in part:
                texts.append(part["text"])
            else:
                texts.append(str(part))
        return "".join(texts)
    return str(content)

# ---------------------------------------------------------------------------
# 2. Researcher Agent
# ---------------------------------------------------------------------------

RESEARCHER_SYSTEM = """
당신은 탁월한 전문성을 보유한 **B2B 시장 동향 및 산업 인텔리전스 리서처 AI** (시장 조사 에이전트)입니다.
귀하의 목표는 설정된 타겟 환경을 기반으로 고품질의 B2B 시장 인사이트 보고서를 조사·편찬하는 것입니다:
- **대상 산업군 (Target Industry)**: {target_industry}
- **구매자 페르소나 (Buyer Persona)**: {buyer_persona}
- **B2B 마케팅 전략 목표 (B2B Strategy Focus)**: {b2b_strategy}

다음 항목들을 명확하고 사실 중심적으로 구조화하여 심층 보고서를 작성해 주십시오:
1. **구매자 페르소나의 핵심 비즈니스 문제(Pain Point) 및 구매 동인 분석**: 해당 산업군 내 타겟 구매자 페르소나가 직면하는 현실적인 기술·비즈니스적 고충, 의사결정을 저해하는 주요 장벽(Objection), 그리고 구매 행동을 촉발하는 핵심 트리거(Purchasing Trigger)를 심층 분석합니다.
2. **글로벌 시장 트렌드 및 규제·환경 컴플라이언스 현황**: 타겟 산업에서 진행 중인 주요 시장 변화 방향, ESG/친환경/지속가능성 관련 글로벌 규제 동향, 그리고 LG화학 등 소재·화학·배터리 분야에서 특히 중요한 green battery 스펙이나 탄소저감 의무화 이슈 등 컴플라이언스 핵심 요소를 정리합니다.
3. **B2B 마케팅 채널 전략 및 핵심 메시징 방향 제안**: 타겟 구매 페르소나에게 가장 효과적으로 도달할 수 있는 최적의 B2B 마케팅 채널 믹스(LinkedIn, 학술/기술 전시회, 기술 백서, 리퍼럴 네트워크 등)와, 이들의 구매 동인에 직접 소구하는 핵심 마케팅 메시지 방향성을 제안합니다.

CRITICAL: 회사명, 특정 제품명, 글로벌 규제/인증 명칭(예: ISCC+, GRS, PCF, TCO, ROI, ESG, EVBattery 등), 전문 기술 스펙 표기(예: NCM, LFP, Tensile Strength, Thermal Stability 등)는 억지로 한글로 번역하지 마시고 영문 고유 원어 표기를 그대로 유지하십시오. 단, 이를 설명하고 분석하는 보고서 전체 본문의 모든 문장 및 서술 내용은 전문적이고 신뢰감 있는 **비즈니스 한국어(완벽한 한국어 비즈니스 문체)**로 풍부하고 상세하게 작성해 주십시오. Markdown 형식으로 구조화하여 출력하십시오.
"""

def researcher_node(state: CampaignState, api_key: str) -> CampaignState:
    _report_stage(state, "researcher")
    client = _get_client(api_key)
    
    prompt = (
        f"# B2B Market Research Request\n"
        f"- Target Industry: {state.get('target_industry', 'Chemicals & Advanced Materials')}\n"
        f"- Buyer Persona / Role: {state.get('buyer_persona', 'Purchasing & Procurement Manager')}\n"
        f"- B2B Campaign Goal: {state.get('b2b_strategy', 'Lead Generation')}\n\n"
        f"Internal Knowledge context (if any):\n{state['knowledge_text'] or '(None)'}"
    )
    
    config_params = {}
    if state.get("google_search"):
        config_params["tools"] = [types.Tool(google_search=types.GoogleSearchRetrieval())]
        
    config = types.GenerateContentConfig(
        system_instruction=RESEARCHER_SYSTEM.format(
            target_industry=state.get('target_industry', 'Chemicals & Advanced Materials'),
            buyer_persona=state.get('buyer_persona', 'Purchasing & Procurement Manager'),
            b2b_strategy=state.get('b2b_strategy', 'Lead Generation')
        ),
        temperature=0.4,
        **config_params
    )
    
    try:
        resp = client.models.generate_content(
            model="gemini-3-flash-preview",
            contents=prompt,
            config=config
        )
        research_data = resp.text
        logger.info("Researcher node done.")
        return {**state, "research_data": research_data, "error": None}
    except Exception as e:
        logger.exception("Researcher error")
        return {**state, "error": str(e)}

# ---------------------------------------------------------------------------
# 2-B. Competitor Marketing Status Agent
# ---------------------------------------------------------------------------

COMPETITOR_SYSTEM = """
당신은 독보적인 전문성을 가진 **B2B 경쟁사 마케팅 인텔리전스 분석가 AI** (경쟁사 분석 에이전트)입니다.
귀하의 목표는 설정된 타겟 환경을 바탕으로 업계 주요 경쟁사들의 마케팅 현황과 채널 전략을 정밀하게 분석하여 보고서를 작성하는 것입니다:
- **대상 산업군 (Target Industry)**: {target_industry}
- **구매자 페르소나 (Buyer Persona)**: {buyer_persona}
- **B2B 마케팅 전략 목표 (B2B Strategy Focus)**: {b2b_strategy}

다음 항목들을 명확하고 상세하게 구조화하여 분석 보고서를 작성해 주십시오:
1. **주요 경쟁사 프로필 및 포지셔닝 분석**: 해당 타겟 시장에서 활동 중인 상위 3개 경쟁사(또는 시장 선도 경쟁사 유형)를 식별하고, 각 기업의 핵심 강점과 시장 내 B2B 포지셔닝을 분석합니다.
2. **경쟁사 마케팅 소구점 및 채널 전략**: 경쟁사들이 주로 사용하는 세일즈 소구 메시지, 활성화된 마케팅 채널(예: LinkedIn, 전문 전시회, 학회, 기술 백서 등), 그리고 이들의 핵심 시각적/브랜드 톤앤매너를 규명합니다.
3. **자사 제품의 차별화 소구 기회 (Strategic Gap & Opportunity)**: 경쟁사들의 소구 포인트 분석을 통해 발견된 시장의 공백이나 약점을 파악하고, 자사 제품이 이를 어떻게 파고들어 차별화된 승리를 거둘 수 있을지에 대한 구체적인 마케팅 침투 전략을 제시합니다.

CRITICAL: 경쟁사 사명(예: LG화학, BASF, CATL 등), 고유 제품명, 특정 기술 키워드 및 비즈니스 전문 용어(예: TCO, ROI, PCF, ISCC+ 등)는 무리하게 억지로 한글로 번역하지 마시고 영문 또는 고유 원어 그대로 표기하십시오. 단, 이를 설명하고 서술하는 본문의 모든 문장과 분석 내용은 전문적이고 세련되며 정중한 **비즈니스 한국어(완벽한 한국어 비즈니스 문체)**로 상세하게 작성해 주셔야 합니다. Markdown 포맷을 사용하여 신뢰도 높은 보고서 형태로 출력하십시오.
"""

def competitor_node(state: CampaignState, api_key: str) -> CampaignState:
    _report_stage(state, "competitor")
    client = _get_client(api_key)
    
    prompt = (
        f"# B2B Competitor Analysis Request\n"
        f"- Target Industry: {state.get('target_industry', 'Chemicals & Advanced Materials')}\n"
        f"- Buyer Persona / Role: {state.get('buyer_persona', 'Purchasing & Procurement Manager')}\n"
        f"- B2B Campaign Goal: {state.get('b2b_strategy', 'Lead Generation')}\n\n"
        f"Internal Knowledge context (if any):\n{state['knowledge_text'] or '(None)'}\n\n"
        f"Market Research Context:\n{state.get('research_data', '(None)')}"
    )
    
    config_params = {}
    if state.get("google_search"):
        config_params["tools"] = [types.Tool(google_search=types.GoogleSearchRetrieval())]
        
    config = types.GenerateContentConfig(
        system_instruction=COMPETITOR_SYSTEM.format(
            target_industry=state.get('target_industry', 'Chemicals & Advanced Materials'),
            buyer_persona=state.get('buyer_persona', 'Purchasing & Procurement Manager'),
            b2b_strategy=state.get('b2b_strategy', 'Lead Generation')
        ),
        temperature=0.4,
        **config_params
    )
    
    try:
        resp = client.models.generate_content(
            model="gemini-3-flash-preview",
            contents=prompt,
            config=config
        )
        competitor_data = resp.text
        logger.info("Competitor node done.")
        return {**state, "competitor_data": competitor_data, "error": None}
    except Exception as e:
        logger.exception("Competitor error")
        return {**state, "error": str(e)}

# ---------------------------------------------------------------------------
# 2-C. Keyword Analysis Agent
# ---------------------------------------------------------------------------

KEYWORD_SYSTEM = """
당신은 글로벌 B2B 검색 엔진 최적화 및 AEO(Answer Engine Optimization) 전략의 최전선에 있는 **B2B 검색·키워드·AEO 전략가 AI** (키워드 분석 에이전트)입니다.
귀하의 목표는 설정된 타겟 환경을 분석하여 검색량이 높고 전환율이 우수한 고관여 핵심 키워드, 타겟 검색 의도, AEO 최적화 전략, 그리고 최적의 소셜 해시태그를 도출하는 것입니다:
- **대상 산업군 (Target Industry)**: {target_industry}
- **구매자 페르소나 (Buyer Persona)**: {buyer_persona}
- **B2B 마케팅 전략 목표 (B2B Strategy Focus)**: {b2b_strategy}

다음 항목들을 실증적인 데이터를 기획하는 것처럼 명확하게 구조화하여 정리해 주십시오:
1. **B2B 고관여 타겟 검색 키워드 (SEO)**: 페르소나의 실질적인 정보 획득 및 구매 의사결정 경로(상업적/트랜잭션 의도)를 반영하여 검색 효율이 우수한 핵심 B2B 키워드 리스트를 제안합니다. (비용 효율, 기술 규격, 품질 인증 등과 연계)
2. **AEO (Answer Engine Optimization) 전략 - AI 검색 시대 대응**: Google AI Overview, ChatGPT Search, Perplexity 등 AI 기반 검색 엔진이 답변을 생성하는 방식에 최적화된 콘텐츠 전략을 분석합니다. AI가 즐겨 인용하는 구조화된 FAQ 형식의 질문-답변 콘텐츠, 전문가 권위(E-E-A-T) 기반 기술 문서 작성 방향, 그리고 AI 검색 Snippet에 포함될 가능성이 높은 명확하고 구체적인 핵심 데이터 포인트(수치, 스펙, 인증 정보)를 제안합니다.
3. **산업군 특화 기술 용어 및 트렌드 검색어**: 대상 산업 분야의 R&D 및 엔지니어, 구매팀이 주로 찾아보는 특화 기술 스펙 명칭과 최근 트렌디하게 급부상하는 비즈니스 핵심 용어들을 도출합니다.
4. **전문 비즈니스 소셜 해시태그 및 SEO/AEO 통합 실행 전략**: 링크드인(LinkedIn) 및 글로벌 B2B 비즈니스 네트워크 채널에서 가시성을 높일 수 있는 전략적 해시태그 조합과, 기존 SEO와 AI 시대의 AEO를 동시에 공략하는 통합 디지털 마케팅 실행 전략을 제시합니다.

CRITICAL: 도출된 핵심 키워드나 제품명, 특정 화학/배터리/IT 기술 규격 단어(예: NCM battery, Tensile Strength, ISCC Plus, Eco-friendly polymer 등)와 소셜 해시태그(예: #EVBattery, #B2BMarketing), AEO 관련 플랫폼명(예: ChatGPT Search, Perplexity, Google AI Overview)은 강제로 억지로 한글로 고쳐 적지 마시고, 영문 고유 표기를 온전히 유지하십시오. 단, 이러한 키워드를 추출한 배경, 타겟별 검색 의도 분석, SEO/AEO 노출 전략 제안 등 모든 분석 서술 문장은 정중하고 깊이 있는 신뢰감을 주는 **비즈니스 한국어**로 완벽하게 상세히 구성되어야 합니다. Markdown 형식으로 깔끔하게 정돈하여 출력해 주십시오.
"""

def keyword_node(state: CampaignState, api_key: str) -> CampaignState:
    _report_stage(state, "keyword")
    client = _get_client(api_key)
    
    prompt = (
        f"# B2B Keyword Analysis Request\n"
        f"- Target Industry: {state.get('target_industry', 'Chemicals & Advanced Materials')}\n"
        f"- Buyer Persona / Role: {state.get('buyer_persona', 'Purchasing & Procurement Manager')}\n"
        f"- B2B Campaign Goal: {state.get('b2b_strategy', 'Lead Generation')}\n\n"
        f"Internal Knowledge context (if any):\n{state['knowledge_text'] or '(None)'}\n\n"
        f"Competitor Strategy context (if any):\n{state.get('competitor_data', '(None)')}"
    )
    
    config_params = {}
    if state.get("google_search"):
        config_params["tools"] = [types.Tool(google_search=types.GoogleSearchRetrieval())]
        
    config = types.GenerateContentConfig(
        system_instruction=KEYWORD_SYSTEM.format(
            target_industry=state.get('target_industry', 'Chemicals & Advanced Materials'),
            buyer_persona=state.get('buyer_persona', 'Purchasing & Procurement Manager'),
            b2b_strategy=state.get('b2b_strategy', 'Lead Generation')
        ),
        temperature=0.4,
        **config_params
    )
    
    try:
        resp = client.models.generate_content(
            model="gemini-3-flash-preview",
            contents=prompt,
            config=config
        )
        keyword_data = resp.text
        logger.info("Keyword node done.")
        return {**state, "keyword_data": keyword_data, "error": None}
    except Exception as e:
        logger.exception("Keyword error")
        return {**state, "error": str(e)}

# ---------------------------------------------------------------------------
# 3. Planner Agent
# ---------------------------------------------------------------------------

PLANNER_SYSTEM = """
You are a **Senior B2B Marketing Strategist AI** (Strategic Planner Agent) specializing in Commercial Excellence.
Your job is to analyze the provided brand knowledge assets, target B2B settings, Market Research report, Competitor Insights, and Keyword Analysis.
Output a highly-strategic **B2B Campaign Strategy Playbook** in markdown.

CRITICAL: Regardless of the language of the input brand knowledge assets, target settings, or user prompts (even if they are in English, Chinese, etc.), you MUST write the entire B2B Campaign Strategy Playbook strictly and completely in natural, professional, high-quality Korean (반드시 모든 내용 및 기획서를 완벽한 한국어로 상세히 작성). Use an authoritative, clear business tone suitable for Korean B2B commercial excellence and marketing teams. DO NOT translate terms literally if they have established professional Korean equivalents.

Your playbook must map the product/brand advantages to the buyer persona's core drivers (e.g. procurement cares about Total Cost of Ownership/ROI; sustainability auditor cares about eco-certification/carbon footprint/regulation; engineer cares about technical reliability/spec sheets).

Output format (strict):
## B2B 캠페인 마케팅 기획서 (Strategy Playbook)
- **B2B 캠페인 목표 (B2B Campaign Objective)**: <한 문장의 핵심 마케팅 전략 목표 - 반드시 완벽한 한국어로 작성>
- **타겟 고객군 및 구매자 페르소나 (Target Audience & Buyer Persona)**: <세부 타겟 기업 설명 및 페르소나가 당면한 현실적인 고충(Pain Point)과 구매 동인 분석 - 반드시 완벽한 한국어로 작성>
- **핵심 가치 드라이버 (Core Value Driver)**: <ROI/비용 절감, ESG/친환경 규제 준수, 독보적인 기술 규격, 또는 파트너십 신뢰 중 핵심 포커스 선정 - 반드시 완벽한 한국어로 작성>
- **핵심 메시지 및 주요 세일즈 포인트 (Core Message & Key Claims)**: <데이터 및 실증 자료에 기반한 고영향력 제품 세일즈 카피 라이팅 - 반드시 완벽한 한국어로 작성>
- **캠페인 톤앤매너 (Tone & Voice)**: <B2B 시장에 신뢰감을 주는 정중하고 기술 중심적인 전문 마케팅 톤앤매너 설계 - 반드시 완벽한 한국어로 작성>
- **브랜드 준수사항 및 규제 제약 요소 (Brand & Regulatory Constraints)**: <LG화학 등 자사 가이드라인 부합 요소 및 준수해야 할 글로벌 B2B 규제 정보 수록 - 반드시 완벽한 한국어로 작성>
- **전략적 광고 카피 라이팅 프레임워크 (Strategic Copywriting Framework)**: <각 구매 주체별 3가지 맞춤형 광고 소구점 구성 전략 설계 - 반드시 완벽한 한국어로 작성>
- **전문 B2B 시각 디자인 가이드라인 (Professional B2B Visual Style)**: <컨셉 이미지 생성을 위한 고품질 테크니컬/연구실/친환경 공장 비주얼 디렉션 제안 - 반드시 완벽한 한국어로 작성>
"""

def planner_node(state: CampaignState, api_key: str) -> CampaignState:
    _report_stage(state, "planner")
    client = _get_client(api_key)
    knowledge = state["knowledge_text"] or "(No additional knowledge assets provided.)"
    
    prompt = (
        f"# B2B Brand Knowledge Assets\n{knowledge}\n\n"
        f"# B2B Market Research Data\n{state.get('research_data', 'No market research data available.')}\n\n"
        f"# B2B Competitor Marketing Insights\n{state.get('competitor_data', 'No competitor marketing insights available.')}\n\n"
        f"# B2B Keyword & Intent Analysis\n{state.get('keyword_data', 'No keyword analysis available.')}\n\n"
        f"# User Prompt / Campaign Request\n{state['user_prompt']}\n\n"
        f"# Target settings:\n"
        f"- Target Industry: {state.get('target_industry', 'Chemicals & Advanced Materials')}\n"
        f"- Buyer Persona: {state.get('buyer_persona', 'Purchasing & Procurement Manager')}\n"
        f"- B2B Campaign Goal: {state.get('b2b_strategy', 'Lead Generation')}\n"
    )
    
    config_params = {}
    if state.get("google_search"):
        config_params["tools"] = [types.Tool(google_search=types.GoogleSearchRetrieval())]
    
    config = types.GenerateContentConfig(
        system_instruction=PLANNER_SYSTEM,
        temperature=0.6,
        **config_params
    )
    
    try:
        resp = client.models.generate_content(
            model="gemini-3-flash-preview",
            contents=prompt,
            config=config
        )
        strategy = resp.text
        logger.info("Planner done.")
        return {**state, "strategy": strategy, "error": None}
    except Exception as e:
        logger.exception("Planner error")
        return {**state, "error": str(e)}

# ---------------------------------------------------------------------------
# 4. Creator Agent (B2B Copywriter)
# ---------------------------------------------------------------------------

CREATOR_SYSTEM = """
You are a **Creative B2B Copywriter AI** (Creator Agent).
Given the B2B Campaign Strategy Playbook, produce the following and NOTHING ELSE:

## Ad Copy Variants
(Provide exactly 3 variants, each clearly labeled as Variant 1, Variant 2, and Variant 3)

You MUST follow this strategic structure for the copies:
- **Variant 1 (Value & ROI Focus)**: Focus on financial performance, total cost of ownership (TCO), scalability, supply chain reliability, and return on investment (ROI). Perfect for procurement managers and C-level.
- **Variant 2 (Technical & Spec Focus)**: Focus on concrete data sheets, chemical/materials/battery specs, engineering stability, rigorous lab testing, and safety standards. Perfect for technical R&D and product engineers.
- **Variant 3 (Relationship & Success Focus)**: Focus on industry leadership, custom co-development partnerships, sustainability/ECO compliance, customer success records, and long-term brand trust. Perfect for commercial and sustainability auditors/directors.

CRITICAL: Regardless of the language of the input strategy playbook, user settings, or any external inputs (even if they are in English or Chinese), you MUST write all the ad copy headlines, bodies, and call-to-actions strictly and completely in highly professional, natural Korean (반드시 모든 카피의 제목, 본문, 행동유도문구를 완벽한 한국어로 작성). The Korean B2B commercial voice must be authoritative, persuasive, and aligned with standard Korean business vocabulary. Do NOT write the copy bodies in English.

Keep the tone and constraints from the strategy playbook strictly. Ensure the copy sounds professional and authoritative.
"""

def creator_node(state: CampaignState, api_key: str) -> CampaignState:
    _report_stage(state, "creator")
    client = _get_client(api_key)
    feedback = ""
    if state.get("review_result", "").startswith("REVISE"):
        feedback = f"\n\n⚠️ Previous attempt was rejected. Critic feedback:\n{state['review_result']}\nPlease revise accordingly."

    prompt = (
        f"# B2B Campaign Strategy Playbook\n{state['strategy']}\n\n"
        f"# Target settings:\n"
        f"- Target Industry: {state.get('target_industry', 'Chemicals & Advanced Materials')}\n"
        f"- Buyer Persona: {state.get('buyer_persona', 'Purchasing & Procurement Manager')}\n"
        f"- B2B Campaign Goal: {state.get('b2b_strategy', 'Lead Generation')}\n"
        f"{feedback}"
    )
    
    config_params = {}
    if state.get("google_search"):
        config_params["tools"] = [types.Tool(google_search=types.GoogleSearchRetrieval())]
    
    config = types.GenerateContentConfig(
        system_instruction=CREATOR_SYSTEM,
        temperature=0.7,
        **config_params
    )

    try:
        resp = client.models.generate_content(
            model="gemini-3-flash-preview",
            contents=prompt,
            config=config
        )
        raw = resp.text
        logger.info(f"Creator raw output length: {len(raw)}")

        copies = []
        
        # Extract ## Ad Copy section if present, else use full raw response
        import re
        if "## Ad Copy" in raw:
            copy_part = raw.split("## Ad Copy", 1)[1]
        else:
            copy_part = raw
            
        # Find all matches of "Variant \d" (allowing surrounding markdown bold markers, bullets, numbering, spaces, and Korean terms like 베리언트, 시안, 버전, 옵션)
        matches = list(re.finditer(r'(?i)(?:\*|-|\d+\.|\s)*(?:Variant|베리언트|시안|버전|옵션)\s*(\d+)', copy_part))
        
        if matches:
            for idx, match in enumerate(matches):
                start = match.start()
                end = matches[idx+1].start() if idx + 1 < len(matches) else len(copy_part)
                var_text = copy_part[start:end].strip()
                if var_text:
                    copies.append(var_text)

        # Fallbacks
        if not copies:
            copies = [
                "Variant 1 (Value/ROI Focus): High efficiency LG Chem B2B materials that reduce TCO and boost productivity.",
                "Variant 2 (Technical/Spec Focus): Advanced grade chemical compounds optimized for thermal stability and high tensile strength.",
                "Variant 3 (Relationship/Trust Focus): Partnering with LG Chem to co-develop sustainable green materials that meet global ECO standards."
            ]

        logger.info(f"Creator done. Copies: {len(copies)}")
        return {**state, "copies": copies, "error": None}
    except Exception as e:
        logger.exception("Creator error")
        return {**state, "error": str(e)}

# ---------------------------------------------------------------------------
# 5. Designer Agent (B2B Visual Director)
# ---------------------------------------------------------------------------

DESIGNER_SYSTEM = """
You are a **B2B Visual Design Director AI** (Designer Agent).
Given the B2B Campaign Strategy Playbook and the generated Ad Copy Variants, design exactly 3 detailed image prompts corresponding to each variant.

Each image prompt must describe a professional B2B visual environment and start with: "Photorealistic marketing image, ..."
Ensure prompts use professional, corporate, and clean B2B aesthetics (e.g. clean technical labs, advanced automatic factories, corporate boardrooms, premium product closeups with microstructures or advanced materials, clean data visualization dashboards).
Avoid generic cartoon, vector, or childish clip art styles.

Provide the prompts in the following markdown section:

## Image Generation Prompts
(Provide exactly 3 detailed image prompts corresponding to each variant.
 Each prompt must start with: "Photorealistic marketing image, ...")
"""

def designer_node(state: CampaignState, api_key: str) -> CampaignState:
    _report_stage(state, "creator")  # designer is part of the creator phase visually
    client = _get_client(api_key)
    feedback = ""
    if state.get("review_result", "").startswith("REVISE"):
        feedback = f"\n\n⚠️ Previous attempt was rejected. Critic feedback:\n{state['review_result']}\nPlease revise accordingly."

    prompt = (
        f"# B2B Campaign Strategy Playbook\n{state['strategy']}\n\n"
        f"# Ad Copy Variants\n" + "\n".join(state.get("copies", [])) + "\n\n"
        f"{feedback}"
    )

    config_params = {}
    config = types.GenerateContentConfig(
        system_instruction=DESIGNER_SYSTEM,
        temperature=0.6,
        **config_params
    )

    try:
        resp = client.models.generate_content(
            model="gemini-3-flash-preview",
            contents=prompt,
            config=config
        )
        raw = resp.text
        logger.info(f"Designer raw output length: {len(raw)}")

        img_prompts = []
        
        current_section = None
        for line in raw.splitlines():
            line_s = line.strip()
            if not line_s: continue
            
            if "## Image Generation" in line_s or "## Image Prompts" in line_s or "Prompt" in line_s:
                current_section = "img"
                continue
            
            if current_section == "img" or line_s.startswith("Photorealistic"):
                if any(line_s.startswith(x) for x in ["Photorealistic", "Prompt", "1.", "2.", "3.", "-"]):
                    img_prompts.append(line_s)
                elif img_prompts:
                    img_prompts[-1] += " " + line_s

        # Clean prompts to ensure they start with Photorealistic...
        cleaned_prompts = []
        for p in img_prompts:
            p_clean = p.strip()
            # Strip prefixes like "Prompt 1: ", "1. ", etc.
            if "Photorealistic" in p_clean:
                idx = p_clean.find("Photorealistic")
                p_clean = p_clean[idx:]
            else:
                p_clean = "Photorealistic marketing image, " + p_clean
            cleaned_prompts.append(p_clean)

        # Fallbacks
        if not cleaned_prompts:
            cleaned_prompts = [
                "Photorealistic marketing image, modern clean corporate office with a large glass table and a data visualization dashboard displaying chemicals production analytics, vibrant blue tone.",
                "Photorealistic marketing image, premium advanced materials closeup, highlighting glossy polymer microstructures inside a clean R&D lab, high precision lighting.",
                "Photorealistic marketing image, professional B2B business partnership shaking hands in a futuristic eco-friendly processing facility, sunny warm daylight, green environment."
            ]

        # Ensure we have exactly 3
        while len(cleaned_prompts) < 3:
            cleaned_prompts.append("Photorealistic marketing image, professional B2B aesthetic, high resolution.")

        logger.info(f"Designer done. Prompts: {len(cleaned_prompts)}")
        return {**state, "image_prompts": cleaned_prompts[:3], "error": None}
    except Exception as e:
        logger.exception("Designer error")
        return {**state, "error": str(e)}

# ---------------------------------------------------------------------------
# 6. Critic Agent
# ---------------------------------------------------------------------------

CRITIC_SYSTEM = """
You are a **B2B Brand Safety & Quality Auditor AI** (Critic Agent).
Evaluate the provided B2B ad copies and image prompts against the B2B Strategy Playbook.

Check for:
1. B2B Tone Alignment: Is the tone authoritative, professional, and data-credible?
2. Brand & Industry compliance constraints.
3. Logical consistency between copy variants (ROI vs Technical vs Trust) and their visual prompts.
4. Factual correctness (check for obvious errors or unsafe B2B claims).

Respond with EXACTLY one of:
- "APPROVE" — if everything is acceptable
- "REVISE: <concise bullet-point feedback on copies/prompts>" — if changes are needed
"""

def critic_node(state: CampaignState, api_key: str) -> CampaignState:
    _report_stage(state, "critic")
    client = _get_client(api_key)
    content = (
        f"# B2B Campaign Strategy Playbook\n{state['strategy']}\n\n"
        f"# Ad Copy Variants\n" + "\n".join(state.get("copies", [])) + "\n\n"
        f"# Image Generation Prompts\n" + "\n".join(state.get("image_prompts", []))
    )
    
    config_params = {}
    if state.get("google_search"):
        config_params["tools"] = [types.Tool(google_search=types.GoogleSearchRetrieval())]

    config = types.GenerateContentConfig(
        system_instruction=CRITIC_SYSTEM,
        temperature=0.4,
        **config_params
    )
    
    try:
        resp = client.models.generate_content(
            model="gemini-3-flash-preview",
            contents=content,
            config=config
        )
        review = resp.text.strip()
        logger.info(f"Critic result: {review[:80]}")
        return {**state, "review_result": review, "error": None}
    except Exception as e:
        logger.exception("Critic error")
        return {**state, "error": str(e), "review_result": "APPROVE"}  # fail-open

# ---------------------------------------------------------------------------
# 7. Routing logic
# ---------------------------------------------------------------------------

MAX_RETRIES = 2

def should_revise(state: CampaignState) -> Literal["revise", "finalize"]:
    review = state.get("review_result", "APPROVE")
    retry = state.get("retry_count", 0)
    if review.startswith("REVISE") and retry < MAX_RETRIES:
        return "revise"
    return "finalize"

def increment_retry(state: CampaignState) -> CampaignState:
    return {**state, "retry_count": state.get("retry_count", 0) + 1}

def finalize_node(state: CampaignState) -> CampaignState:
    return {
        **state,
        "final_copies": state.get("copies", []),
        "final_image_prompts": state.get("image_prompts", []),
    }

# ---------------------------------------------------------------------------
# 8. Build Graph
# ---------------------------------------------------------------------------

def build_campaign_graph(api_key: str):
    graph = StateGraph(CampaignState)

    # Bind api_key to nodes via closures
    graph.add_node("researcher",lambda s: researcher_node(s, api_key))
    graph.add_node("competitor",lambda s: competitor_node(s, api_key))
    graph.add_node("keyword",   lambda s: keyword_node(s, api_key))
    graph.add_node("planner",   lambda s: planner_node(s, api_key))
    graph.add_node("creator",   lambda s: creator_node(s, api_key))
    graph.add_node("designer",  lambda s: designer_node(s, api_key))
    graph.add_node("critic",    lambda s: critic_node(s, api_key))
    graph.add_node("retry",     increment_retry)
    graph.add_node("finalize",  finalize_node)

    graph.set_entry_point("researcher")
    graph.add_edge("researcher", "competitor")
    graph.add_edge("competitor", "keyword")
    graph.add_edge("keyword", "planner")
    graph.add_edge("planner", "creator")
    graph.add_edge("creator", "designer")
    graph.add_edge("designer", "critic")
    graph.add_conditional_edges(
        "critic",
        should_revise,
        {"revise": "retry", "finalize": "finalize"}
    )
    graph.add_edge("retry", "creator")
    graph.add_edge("finalize", END)

    return graph.compile()

# ---------------------------------------------------------------------------
# 9. File Parsers
# ---------------------------------------------------------------------------

def extract_text_from_pdf(file_bytes: bytes) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(file_bytes))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception as e:
        logger.error(f"PDF parse error: {e}")
        return ""

def extract_text_from_docx(file_bytes: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        logger.error(f"DOCX parse error: {e}")
        return ""

def extract_text_from_pptx(file_bytes: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
        return "\n".join(texts)
    except Exception as e:
        logger.error(f"PPTX parse error: {e}")
        return ""

def extract_text_from_txt(file_bytes: bytes) -> str:
    try:
        return file_bytes.decode("utf-8", errors="ignore")
    except Exception as e:
        logger.error(f"TXT parse error: {e}")
        return ""

def parse_file(filename: str, file_bytes: bytes) -> str:
    ext = filename.lower().rsplit(".", 1)[-1]
    if ext == "pdf":
        return extract_text_from_pdf(file_bytes)
    elif ext == "docx":
        return extract_text_from_docx(file_bytes)
    elif ext in ("pptx", "ppt"):
        return extract_text_from_pptx(file_bytes)
    elif ext == "txt":
        return extract_text_from_txt(file_bytes)
    else:
        logger.warning(f"Unsupported file type: {ext}")
        return f"(Unsupported file type: {filename})"
